"""
report_generator.py - HarvestBridge quarterly report generator
Blue accent theme
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io, os, tempfile
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

# Colour palette 
BLUE        = RGBColor(0x1A, 0x5C, 0xA8)   # primary - section labels, underlines, icons
DARK_BLUE   = RGBColor(0x0D, 0x3A, 0x6E)   # header bar background
LIGHT_BLUE  = RGBColor(0xE8, 0xF0, 0xFA)   # subtle tint if needed
GOLD        = RGBColor(0xC9, 0xA8, 0x3C)   # EBITDA bars
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MID_GRAY    = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY   = RGBColor(0x22, 0x22, 0x22)
BORDER      = RGBColor(0xCC, 0xCC, 0xCC)
GREEN       = RGBColor(0x1A, 0x7A, 0x3C)   # all positive/growth indicators
RED         = RGBColor(0xCC, 0x22, 0x22)   # negative indicators
BLACK       = RGBColor(0x00, 0x00, 0x00)
BLUE_LIGHT  = RGBColor(0xB8, 0xCC, 0xE8)   # header subtitle text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# Drawing helpers 

def rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def tb(slide, text, x, y, w, h, size=9, bold=False, color=None,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color or DARK_GRAY
    return box


def section_label(slide, text, x, y, w):
    """Blue section label with underline — matching sample style"""
    tb(slide, text, x, y, w, Inches(0.22), size=8, bold=True, color=BLUE)
    rect(slide, x, y + Inches(0.22), w, Inches(0.016), fill=BLUE)


# Formatters 

def fusd(v):
    if v is None: return "N/A"
    if abs(v) >= 1000: return f"${v/1000:.2f}M"
    return f"${v:.0f}K"

def fpct(v):
    if v is None: return "N/A"
    return f"{v*100:.1f}%"

def fqoq(v):
    if v is None: return ""
    sign = "+" if v >= 0 else ""
    return f"{sign}{v*100:.1f}%"

def fnum(v):
    if v is None: return "N/A"
    if isinstance(v, str): return v
    if v >= 1000: return f"{v/1000:.1f}K"
    return str(int(v))

def fmt_kpi_pct(v):
    if v is None: return "N/A"
    f = float(v)
    return f"{f*100:.0f}%" if f <= 1.0 else f"{f:.0f}%"


# Icons (matplotlib, blue theme) 

def _save_fig(fig, dpi=120):
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    plt.savefig(tmp.name, dpi=dpi, bbox_inches="tight", transparent=True)
    plt.close()
    return tmp.name


def icon_loans():
    fig, ax = plt.subplots(figsize=(0.6, 0.6))
    fig.patch.set_alpha(0); ax.set_axis_off()
    ax.set_xlim(0,1); ax.set_ylim(0,1)
    col = "#1A5CA8"
    ax.plot([0.5,0.5],[0.1,0.85], color=col, lw=2.5)
    for y in [0.3,0.45,0.6,0.75]:
        ax.plot([0.5,0.3],[y,y+0.1], color=col, lw=1.8)
        ax.plot([0.5,0.7],[y,y+0.1], color=col, lw=1.8)
    ax.add_patch(mpatches.Ellipse((0.5,0.88),0.12,0.1,color=col))
    return _save_fig(fig)


def icon_women():
    fig, ax = plt.subplots(figsize=(0.6, 0.6))
    fig.patch.set_alpha(0); ax.set_axis_off()
    ax.set_xlim(0,1); ax.set_ylim(0,1)
    col = "#1A5CA8"
    ax.add_patch(plt.Circle((0.5,0.78),0.14,color=col))
    ax.fill([0.25,0.75,0.5],[0.28,0.28,0.62],color=col)
    ax.plot([0.3,0.22],[0.55,0.3],color=col,lw=2)
    ax.plot([0.7,0.78],[0.55,0.3],color=col,lw=2)
    return _save_fig(fig)


def icon_tractor():
    fig, ax = plt.subplots(figsize=(0.6, 0.6))
    fig.patch.set_alpha(0); ax.set_axis_off()
    ax.set_xlim(0,1); ax.set_ylim(0,1)
    col = "#1A5CA8"
    ax.add_patch(plt.Rectangle((0.2,0.35),0.55,0.32,color=col,ec=col))
    ax.add_patch(plt.Rectangle((0.48,0.55),0.28,0.22,color=col,ec=col))
    ax.add_patch(plt.Circle((0.32,0.3),0.14,color="#555",ec="#555"))
    ax.add_patch(plt.Circle((0.32,0.3),0.07,color=col))
    ax.add_patch(plt.Circle((0.67,0.32),0.1,color="#555",ec="#555"))
    ax.add_patch(plt.Circle((0.67,0.32),0.05,color=col))
    return _save_fig(fig)


def icon_revenue():
    fig, ax = plt.subplots(figsize=(0.6, 0.6))
    fig.patch.set_alpha(0); ax.set_axis_off()
    ax.set_xlim(0,1); ax.set_ylim(0,1)
    col = "#1A5CA8"
    ax.add_patch(plt.Circle((0.5,0.5),0.42,color=col,ec=col))
    ax.text(0.5,0.5,"$",ha="center",va="center",
            fontsize=22,fontweight="bold",color="white")
    return _save_fig(fig)


ICON_FUNCS = [icon_loans, icon_women, icon_tractor, icon_revenue]


# Logo 

def make_logo():
    fig, ax = plt.subplots(figsize=(1.5, 0.38))
    fig.patch.set_facecolor("#0D3A6E"); ax.set_axis_off()
    ax.set_xlim(0,1); ax.set_ylim(0,1)
    ax.text(0.5,0.62,"HarvestBridge",ha="center",va="center",
            fontsize=9,fontweight="bold",color="white")
    ax.text(0.5,0.25,"Agri-Fintech",ha="center",va="center",
            fontsize=6,color="#B8CCE8")
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    plt.savefig(tmp.name, dpi=120, bbox_inches="tight", facecolor="#0D3A6E")
    plt.close()
    return tmp.name


# Bar chart 

def make_chart(quarters, fin_hist, new_label, new_rev, new_ebitda):
    labels  = quarters + [new_label]
    revenue = [fin_hist[q]["total_revenue"]/1000 for q in quarters] + [new_rev/1000]
    ebitda  = [fin_hist[q]["ebitda"]/1000        for q in quarters] + [new_ebitda/1000]

    x  = np.arange(len(labels))
    bw = 0.35

    fig, ax = plt.subplots(figsize=(5.8, 2.6))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    b1 = ax.bar(x - bw/2 - 0.03, revenue, width=bw,
                color="#1A5CA8", label="Total Revenue", zorder=3)
    b2 = ax.bar(x + bw/2 + 0.03, ebitda,  width=bw,
                color="#C9A83C", label="EBITDA",        zorder=3)

    for bar in list(b1):
        h = bar.get_height()
        ax.text(bar.get_x()+bar.get_width()/2, h+0.02,
                f"${h:.2f}M", ha="center", va="bottom",
                fontsize=5.8, color="#1A5CA8", fontweight="bold")
    for bar in list(b2):
        h = bar.get_height()
        label = f"${h*1000:.0f}K" if abs(h) < 0.5 else f"${h:.2f}M"
        ypos  = h - 0.07 if h < 0 else h + 0.02
        va    = "top" if h < 0 else "bottom"
        ax.text(bar.get_x()+bar.get_width()/2, ypos,
                label, ha="center", va=va,
                fontsize=5.8, color="#9A7820", fontweight="bold")

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=7, color="#444")
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"${v:.1f}M"))
    ax.tick_params(axis="y", labelsize=6.5, colors="#666")
    ax.set_ylim(min(min(ebitda)*1.5, -0.15), max(revenue)*1.28)
    ax.set_title("Key Financial Metrics (USD '000)", fontsize=7.5,
                 color="#1A5CA8", fontweight="bold", loc="left", pad=4)
    for sp in ["top","right"]: ax.spines[sp].set_visible(False)
    for sp in ["left","bottom"]: ax.spines[sp].set_color("#CCCCCC")
    ax.grid(axis="y", color="#EEEEEE", lw=0.6, zorder=0)
    ax.legend(fontsize=6.5, loc="upper left", framealpha=0)

    plt.tight_layout(pad=0.5)
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    plt.savefig(tmp.name, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close()
    return tmp.name


# Main 

def generate_report(financials_calc, kpis, quarter_label, prior_label,
                    company_info, historical_financials, historical_quarters):

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tmp_files = []

    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)

    # HEADER 
    hdr_h = Inches(0.44)
    rect(slide, 0, 0, SLIDE_W, hdr_h, fill=DARK_BLUE)
    tb(slide,
       f"HarvestBridge  –  {company_info['country']}  –  {company_info['sector']}  –  {company_info['stake']} of Fund value",
       Inches(0.18), Inches(0.08), Inches(10), Inches(0.32),
       size=14, bold=True, color=WHITE)

    # COLUMN SETUP 
    col_y   = Inches(0.52)
    left_x  = Inches(0.18)
    left_w  = Inches(6.28)
    right_x = Inches(6.72)
    right_w = Inches(6.44)

    # LEFT  RAG indicators 
    rag_items = [("Runway","●●●●"),("M&G","●●●●"),("EBITDA+","●●●"),("Exit readiness","●●●●")]
    rag_h  = Inches(0.30)
    box_w  = Inches(1.44)
    gap    = Inches(0.06)
    rag_y  = col_y + Inches(0.02)

    # 2x2 grid to the right, next to logo
    rag_box_w = Inches(1.3)
    rag_box_h = Inches(0.20)
    rag_col1  = Inches(9.2)
    rag_col2  = rag_col1 + rag_box_w + Inches(0.05)
    positions = [
        (rag_col1, rag_y),
        (rag_col2, rag_y),
        (rag_col1, rag_y + rag_box_h + Inches(0.04)),
        (rag_col2, rag_y + rag_box_h + Inches(0.04)),
    ]
    for (bx, by), (lbl, dots) in zip(positions, rag_items):
        rect(slide, bx, by, rag_box_w, rag_box_h, fill=WHITE, line=BLUE)
        tb(slide, f"{lbl}  {dots}", bx+Inches(0.07), by+Inches(0.03),
           rag_box_w-Inches(0.09), Inches(0.15),
           size=6, bold=False, color=DARK_GRAY)

    # RIGHT - Logo (top right, same row as RAG)
    logo_path = make_logo()
    tmp_files.append(logo_path)
    slide.shapes.add_picture(logo_path,
        Inches(11.85), rag_y, Inches(1.28), Inches(0.44))

    # Column content starts below RAG row
    col_y = rag_y + rag_h + Inches(0.12)

    # Vertical divider
    rect(slide, Inches(6.62), col_y, Inches(0.014), Inches(6.6), fill=LIGHT_GRAY)

    # LEFT - Description
    section_label(slide, "Description of the Business", left_x, col_y, left_w)
    tb(slide, company_info.get("description",""),
       left_x, col_y+Inches(0.28), left_w, Inches(0.92),
       size=7.8, color=DARK_GRAY, wrap=True)

    # LEFT - Company Impact 
    imp_y = col_y + Inches(1.26)
    section_label(slide, "Company Impact", left_x, imp_y, left_w)
    imp_y += Inches(0.28)

    female_pct = 0
    if kpis.get("fte_headcount") and kpis.get("female_fte"):
        female_pct = round(float(kpis["female_fte"])/float(kpis["fte_headcount"])*100,1)

    wb_pct    = kpis.get("women_borrowers_pct")
    women_str = f"{wb_pct:.0f}%" if wb_pct else "58%"
    rev_qoq   = fqoq(financials_calc.get("revenue_qoq"))

    cards = [
        (0, fnum(kpis.get("loans_disbursed")),      "Loans disbursed",          rev_qoq),
        (1, women_str,                               "Women borrowers",           None),
        (2, fnum(kpis.get("active_saas_users")),    "Active farmers",            None),
        (3, fusd(financials_calc["total_revenue"]),  f"{quarter_label} revenue", rev_qoq),
    ]

    # Card dimensions - taller to give breathing room
    cw   = (left_w - Inches(0.06)) / 4
    ch   = Inches(1.12)
    cgap = Inches(0.02)

    for idx, val, lbl, qval in cards:
        cx = left_x + idx*(cw+cgap)

        # Card - white background, light border, no shadow
        rect(slide, cx, imp_y, cw, ch, fill=WHITE, line=BORDER)

        # Icon - top left
        icon_path = ICON_FUNCS[idx]()
        tmp_files.append(icon_path)
        slide.shapes.add_picture(icon_path,
            cx+Inches(0.07), imp_y+Inches(0.06),
            Inches(0.30), Inches(0.30))

        # Value - large, bold
        tb(slide, val,
           cx+Inches(0.07), imp_y+Inches(0.38),
           cw-Inches(0.1), Inches(0.3),
           size=13, bold=True, color=BLACK)

        # Label
        tb(slide, lbl,
           cx+Inches(0.07), imp_y+Inches(0.67),
           cw-Inches(0.1), Inches(0.16),
           size=6.5, color=MID_GRAY)

    # Quarter to quarter row 
    qtq_y = imp_y + ch + Inches(0.02)
    qtq_h = Inches(0.36)
    rect(slide, left_x, qtq_y, left_w, qtq_h, fill=WHITE)

    # Label on its own line
    tb(slide, "Quarter to quarter",
       left_x+Inches(0.07), qtq_y+Inches(0.02),
       left_w, Inches(0.14), size=6.5, color=MID_GRAY)

    # Values aligned to card columns on the line below
    for idx, val, lbl, qval in cards:
        cx      = left_x + idx*(cw+cgap)
        display = qval if qval else "NA"
        dcol    = (GREEN if qval.startswith("+") else RED) if qval else MID_GRAY
        tb(slide, display,
           cx+Inches(0.05), qtq_y+Inches(0.18),
           cw-Inches(0.08), Inches(0.16),
           size=8, bold=bool(qval), color=dcol,
           align=PP_ALIGN.CENTER)

    rect(slide, left_x, qtq_y+qtq_h, left_w, Inches(0.012), fill=LIGHT_GRAY)

    # Growth since investment row
    gsi_y = qtq_y + qtq_h + Inches(0.012)
    gsi_h = Inches(0.36)
    rect(slide, left_x, gsi_y, left_w, gsi_h, fill=OFF_WHITE)

   # Label on its own line
    tb(slide, "Growth since investment",
       left_x+Inches(0.07), gsi_y+Inches(0.02),
       left_w, Inches(0.14), size=6.5, color=MID_GRAY)

    # NA values aligned to card columns on the line below
    for i in range(4):
        cx = left_x + i*(cw+cgap)
        tb(slide, "NA",
           cx+Inches(0.05), gsi_y+Inches(0.18),
           cw-Inches(0.08), Inches(0.14),
           size=8, color=MID_GRAY, align=PP_ALIGN.CENTER)

    qu_y = gsi_y + gsi_h + Inches(0.12)
    section_label(slide, "Quarterly Update", left_x, qu_y, left_w)

    update_text = company_info.get("quarterly_update", "")
    tb(slide, update_text,
       left_x, qu_y+Inches(0.28), left_w, Inches(2.1),
       size=7.8, color=DARK_GRAY, wrap=True)

    # Footer
    tb(slide, "1", Inches(0.15), Inches(7.26), Inches(0.25), Inches(0.18),
       size=7, color=MID_GRAY)
    tb(slide,
       f"Alitheia Capital, Quarterly Report: {quarter_label}  |  Private and confidential",
       Inches(0.44), Inches(7.26), Inches(6.2), Inches(0.18),
       size=7, color=MID_GRAY)

    # RIGHT — Valuation section
    section_label(slide, "Investment Manager's Update on Valuation and Exit",
                  right_x, col_y, right_w)

    val_text = company_info.get("valuation_update",
        f"HarvestBridge is currently valued at USD 33.00M (post-money Series B, Q2 2025). "
        f"Alitheia invested USD 3.50M as Series B lead, acquiring a {company_info['stake']}. "
        f"NAV is marked at USD 4.20M (1.2x cost), reflecting first-quarter profitability and "
        f"strong EBITDA trajectory. No active exit process at this stage; the investment is "
        f"expected to be held for 4 - 5 years with a target exit via strategic sale or secondary.")
    tb(slide, val_text,
       right_x, col_y+Inches(0.28), right_w, Inches(0.95),
       size=7.8, color=DARK_GRAY, wrap=True)

    # RIGHT - Bar chart 
    chart_y = col_y + Inches(1.28)
    chart_h = Inches(2.55)
    chart_path = make_chart(
        historical_quarters, historical_financials,
        quarter_label,
        financials_calc["total_revenue"],
        financials_calc["ebitda"]
    )
    tmp_files.append(chart_path)
    slide.shapes.add_picture(chart_path, right_x, chart_y, right_w, chart_h)

    # RIGHT - Outlook 
    out_y = chart_y + chart_h + Inches(0.08)
    section_label(slide, "Investment Manager's Outlook and Value Addition",
                  right_x, out_y, right_w)

    outlook_text = company_info.get("outlook",
        f"In {quarter_label}, HarvestBridge delivered net income of "
        f"{fusd(financials_calc['net_income'])} and EBITDA of {fusd(financials_calc['ebitda'])} "
        f"({fpct(financials_calc['ebitda_margin'])} margin), validating the operating model's "
        f"scalability. The underlying drivers — improving gross margins "
        f"({fpct(financials_calc['gross_margin'])}), disciplined opex management - are durable.\n\n"
        f"Management is focused on scaling loan disbursements toward 200,000 annually by FY2027, "
        f"deepening Kenya penetration, and extending the SaaS platform. Alitheia continues to "
        f"support gender lens reporting, board-level governance, and the company's path to a "
        f"Series C raise in H2 2027.")
    tb(slide, outlook_text,
       right_x, out_y+Inches(0.28), right_w, Inches(2.1),
       size=7.8, color=DARK_GRAY, wrap=True)

    # Cleanup
    for p in tmp_files:
        try: os.unlink(p)
        except: pass

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
